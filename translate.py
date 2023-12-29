from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from translators import get_translator


def translate_pptx(
    input_path, output_path,
    source_language, target_language,
    translator="google"
):
    """
    Translate all texts in a PowerPoint file, including:
    1. Text in text boxes
    2. Text in tables
    3. Text in images (if possible)
    """
    translator = get_translator(source_language, target_language, translator)
    ppt = Presentation(input_path)

    for slide in ppt.slides:
        translate_slide(slide, translator)
        ppt.save(output_path)


def translate_slide(slide, translator):
    """
    Translate all texts in a slide inplace, including:
    1. Text in text boxes
    2. Text in tables
    3. Text in grouped shapes
    """
    text_boxes = []
    tables = []
    groups = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            text_boxes.append(shape)
        elif shape.has_table and shape.table:
            tables.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            groups.append(shape)

    # translate text boxes
    _translate_text_boxes(text_boxes, translator)
    _translate_tables(tables, translator)
    _translate_groups(groups, translator)


def _translate_text_boxes(targets, translator):
    """
    Translate all texts in text boxes inplace.
    """
    texts = [shape.text for shape in targets]
    translated_text = translator.translate(texts)

    # Replace the original text with the translated text
    for shape, translated_text in zip(targets, translated_text):
        shape.text = translated_text


def _translate_tables(targets, translator):
    """
    Translate all texts in tables inplace.
    """
    # translate texts table by table
    for shape in targets:
        texts = []
        for row in shape.table.rows:
            for cell in row.cells:
                texts.append(cell.text)

        translated_text = translator.translate(texts)

        # pop translated text one by one
        for row in shape.table.rows:
            for cell in row.cells:
                cell.text = translated_text.pop(0)


def _translate_groups(targets, translator):
    """
    Translate all texts in grouped shapes inplace recursively.
    """
    return [translate_slide(shape, translator) for shape in targets]
